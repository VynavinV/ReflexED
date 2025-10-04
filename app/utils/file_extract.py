"""
Helpers to extract text from various document formats for assignments.
"""
from typing import Optional
import os


def extract_text(path: str) -> Optional[str]:
    print(f"📖 Extracting text from: {path}")
    if not path or not os.path.exists(path):
        print(f"⚠️ File not found: {path}")
        return None
    ext = os.path.splitext(path)[1].lower()
    print(f"📄 File type detected: {ext}")
    try:
        if ext in ('.txt', '.md'):
            print("📄 Reading text file...")
            with open(path, 'r', encoding='utf-8', errors='ignore') as f:
                content = f.read()
            print(f"✅ Text file read: {len(content)} characters")
            return content
        elif ext == '.pdf':
            print("📄 Processing PDF file...")
            from pypdf import PdfReader
            reader = PdfReader(path)
            print(f"📄 PDF has {len(reader.pages)} pages")
            out = []
            for i, page in enumerate(reader.pages):
                text = page.extract_text() or ''
                print(f"📄 Page {i+1}: {len(text)} characters")
                out.append(text)
            result = '\n'.join(out).strip() or None
            print(f"✅ PDF processed: {len(result) if result else 0} total characters")
            return result
        elif ext in ('.docx',):
            print("📄 Processing Word document...")
            import docx
            doc = docx.Document(path)
            content = '\n'.join([p.text for p in doc.paragraphs]).strip() or None
            print(f"✅ Word document processed: {len(content) if content else 0} characters")
            return content
        elif ext in ('.pptx',):
            print("📄 Processing PowerPoint presentation...")
            from pptx import Presentation
            prs = Presentation(path)
            out = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, 'text'):
                        out.append(shape.text)
            content = '\n'.join(out).strip() or None
            print(f"✅ PowerPoint processed: {len(content) if content else 0} characters")
            return content
        else:
            print("📄 Unknown file type, attempting text read...")
            # Fallback try text read
            with open(path, 'r', encoding='utf-8', errors='ignore') as f:
                content = f.read()
            print(f"✅ Fallback text read: {len(content)} characters")
            return content
    except Exception as e:
        print(f"❌ Text extraction failed: {e}")
        return None
